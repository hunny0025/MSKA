"""
PPTX text extractor.
"""
from pipeline.extractors import ExtractionError


def extract(file_path: str) -> str:
    """Extract text from a .pptx file from all slide shapes. Raises ExtractionError on failure."""
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        raise ExtractionError("python-pptx is not installed")

    try:
        prs = Presentation(file_path)
        parts = []
        for slide_num, slide in enumerate(prs.slides, start=1):
            parts.append(f"[Slide {slide_num}]")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text)
        result = "\n".join(parts).strip()
        if not result:
            raise ExtractionError(f"PPTX produced empty text: {file_path}")
        return result
    except ExtractionError:
        raise
    except Exception as exc:
        raise ExtractionError(f"PPTX extraction failed for {file_path}: {exc}") from exc
