"""
Dummy Data Seeding Script — Prompt 9A

Creates an isolated SQLite database with:
- 5 departments (QA, Production, HR, Engineering, Supply Chain)
- 5 roles (employee, department_lead, project_admin, platform_admin, auditor)
- 6 users covering every role
- 5 projects across departments
- 6 synthetic documents (PDF/DOCX/XLSX/PPTX/TXT/CSV)

Documents are ingested through the REAL production pipeline
(parsers → PII scanner → classifier → chunker → embedder → FAISS indexer).

One document deliberately contains PII (to test quarantine).
One document is deliberately irrelevant (to test abstention).
Documents span classification levels (public/internal/confidential/restricted).

Usage:
    python scripts/seed_dummy_data.py
    # Or imported as a module by test fixtures
"""

import asyncio
import io
import os
import sys
import uuid
import tempfile
import shutil

import bcrypt

# Ensure the backend package is importable
BACKEND_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "backend")
sys.path.insert(0, BACKEND_DIR)

from sqlalchemy.ext.asyncio import create_async_engine, async_sessionmaker, AsyncSession
from sqlalchemy import select, insert

from core.database import Base
from core.permissions import Roles
from models.user import User
from models.role import Role
from models.department import Department
from models.project import Project
from models.project import project_users
from models.document import Document


# ═══════════════════════════════════════════════════════════════════
# Synthetic Document Content
# ═══════════════════════════════════════════════════════════════════

DOC_CONTENTS = {
    "sop_welding_robot.txt": {
        "text": (
            "Maruti Suzuki Internal Use Only\n\n"
            "STANDARD OPERATING PROCEDURE — WELDING ROBOT CALIBRATION\n\n"
            "SOP Number: QA-WR-2024-003\n"
            "Department: Quality Assurance\n"
            "Revision: 3.1\n\n"
            "1. PURPOSE\n"
            "This SOP defines the step-by-step calibration procedure for Fanuc R-2000iC "
            "welding robots installed on Line 4 at the Gurugram plant. Correct calibration "
            "ensures weld seam integrity within ISO 3834 tolerances.\n\n"
            "2. SCOPE\n"
            "Applies to all welding robots on Body Shop Line 4 (stations WS-401 through WS-412).\n\n"
            "3. SAFETY PRECAUTIONS\n"
            "- Engage LOTO (Lock Out / Tag Out) before entering the robot cell.\n"
            "- Wear heat-resistant gloves, face shield, and steel-toe boots.\n"
            "- Verify teach pendant E-stop functionality before power-on.\n\n"
            "4. CALIBRATION STEPS\n"
            "Step 1: Power down the controller via the main disconnect switch.\n"
            "Step 2: Attach the Renishaw Ballbar QC20-W to the robot flange.\n"
            "Step 3: Run the automated calibration cycle (Program CAL_LINE4_V3).\n"
            "Step 4: Record deviation values in the MES system under ticket WR-CAL.\n"
            "Step 5: If deviation exceeds ±0.15mm, escalate to Maintenance Lead.\n\n"
            "5. ACCEPTANCE CRITERIA\n"
            "All six axes must read within ±0.10mm of nominal.\n"
            "TCP (Tool Center Point) drift must not exceed 0.05mm over 8 hours.\n\n"
            "6. RECORDS\n"
            "Calibration certificates stored in QMS portal under Documents > Welding.\n"
        ),
        "classification": "internal",
        "format": "txt",
    },

    "supplier_audit_report.csv": {
        "text": (
            "Supplier Name,Part Number,Defect Rate (%),Audit Date,Classification\n"
            "Bharat Forge,BF-CRK-7721,0.12,2024-06-15,Confidential\n"
            "Minda Industries,MI-HDL-3305,0.45,2024-06-18,Confidential\n"
            "Sona BLW,SB-GEAR-9901,0.03,2024-07-01,Confidential\n"
            "Lumax Auto,LA-LAMP-5520,1.20,2024-07-05,Confidential\n"
            "Subros Ltd,SUB-AC-8810,0.08,2024-07-10,Confidential\n"
            "Rico Auto,RA-DIE-4402,0.67,2024-07-12,Confidential\n\n"
            "NOTE: This supplier quality data is classified CONFIDENTIAL.\n"
            "Distribution restricted to Production and QA leads only.\n"
        ),
        "classification": "confidential",
        "format": "csv",
    },

    "hr_employee_roster.docx": {
        "text": (
            "HUMAN RESOURCES — EMPLOYEE ROSTER (CONFIDENTIAL)\n\n"
            "Employee Name: Rajesh Kumar\n"
            "Aadhaar Number: 1234 5678 9012\n"
            "PAN Card: ABCDE1234F\n"
            "Phone: +91 9876543210\n"
            "Email: rajesh.kumar@marutisuzuki.com\n"
            "CTC: 12,50,000 INR per annum\n"
            "Department: HR\n\n"
            "Employee Name: Priya Sharma\n"
            "Aadhaar Number: 9876 5432 1098\n"
            "PAN Card: FGHIJ5678K\n"
            "Phone: +91 8765432109\n"
            "Email: priya.sharma@marutisuzuki.com\n"
            "Salary: 8,75,000 INR per annum\n"
            "Department: HR\n"
        ),
        "classification": "confidential",
        "format": "docx",
        "has_pii": True,
    },

    "engine_spec_v2.pdf": {
        "text": (
            "RESTRICTED — TRADE SECRET\n\n"
            "ENGINE SPECIFICATION DOCUMENT — K15C SMART HYBRID\n"
            "Document ID: ENG-SPEC-K15C-V2\n\n"
            "1. ENGINE OVERVIEW\n"
            "The K15C 1.5L Dual VVT petrol engine with integrated starter generator (ISG) "
            "delivers 103 PS at 6000 rpm and 138 Nm torque at 4400 rpm. The smart hybrid "
            "system provides torque assist, regenerative braking, and idle start-stop.\n\n"
            "2. KEY SPECIFICATIONS\n"
            "Displacement: 1462 cc\n"
            "Bore x Stroke: 74.0 mm x 85.0 mm\n"
            "Compression Ratio: 12.5:1\n"
            "Fuel System: Multi-point fuel injection (MPFI)\n"
            "Emission Standard: BS-VI Phase 2\n\n"
            "3. PROPRIETARY COATING\n"
            "Cylinder bore uses plasma-transferred wire-arc coating (PTWA) — "
            "this is a proprietary Suzuki process with restricted distribution.\n"
            "Trade secret: coating thickness tolerance is 150±5 microns.\n\n"
            "4. PERFORMANCE CURVES\n"
            "See Appendix A for full torque and power curves.\n"
        ),
        "classification": "restricted",
        "format": "pdf",
    },

    "safety_training_deck.pptx": {
        "text": (
            "MARUTI SUZUKI — PLANT FLOOR SAFETY TRAINING\n\n"
            "Slide 1: Welcome to Safety Orientation\n"
            "All new employees must complete this training before entering the shop floor.\n\n"
            "Slide 2: Personal Protective Equipment (PPE)\n"
            "- Hard hat (mandatory in all production zones)\n"
            "- Safety goggles (mandatory near grinding/welding)\n"
            "- Steel-toe boots (mandatory everywhere)\n"
            "- High-visibility vest (mandatory in logistics zones)\n\n"
            "Slide 3: Emergency Procedures\n"
            "- Fire alarm: evacuate via nearest green exit sign\n"
            "- Chemical spill: call EHS hotline 1800-XXX-XXXX\n"
            "- Injury: report to nearest first-aid station\n\n"
            "Slide 4: Housekeeping (5S)\n"
            "Sort, Set in Order, Shine, Standardize, Sustain.\n"
            "Keep walkways clear. No tools left on the floor.\n"
        ),
        "classification": "public",
        "format": "pptx",
    },

    "random_recipe.xlsx": {
        "text": (
            "RECIPE: BUTTER CHICKEN (Murgh Makhani)\n\n"
            "Ingredients:\n"
            "500g boneless chicken thighs\n"
            "1 cup tomato puree\n"
            "2 tbsp butter\n"
            "1 cup heavy cream\n"
            "1 tbsp garam masala\n"
            "1 tsp turmeric powder\n"
            "Salt to taste\n\n"
            "Instructions:\n"
            "1. Marinate chicken in yogurt and spices for 2 hours.\n"
            "2. Grill or pan-fry chicken until charred.\n"
            "3. Simmer tomato puree with butter and cream.\n"
            "4. Add chicken to sauce and cook for 15 minutes.\n"
            "5. Garnish with fresh cream and kasuri methi.\n\n"
            "This has absolutely nothing to do with automotive manufacturing.\n"
        ),
        "classification": "internal",
        "format": "xlsx",
    },
}


# ═══════════════════════════════════════════════════════════════════
# Binary File Generators
# ═══════════════════════════════════════════════════════════════════

def _generate_txt_bytes(text: str) -> bytes:
    return text.encode("utf-8")


def _generate_csv_bytes(text: str) -> bytes:
    return text.encode("utf-8")


def _generate_docx_bytes(text: str) -> bytes:
    """Generate a real .docx file from text content."""
    from docx import Document as DocxDocument
    doc = DocxDocument()
    for para in text.split("\n"):
        if para.strip():
            doc.add_paragraph(para)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _generate_pdf_bytes(text: str) -> bytes:
    """Generate a real .pdf file from text content."""
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfgen import canvas as pdf_canvas
    buf = io.BytesIO()
    c = pdf_canvas.Canvas(buf, pagesize=A4)
    width, height = A4
    y = height - 50
    for line in text.split("\n"):
        if y < 50:
            c.showPage()
            y = height - 50
        c.drawString(50, y, line[:100])  # Truncate long lines for PDF layout
        y -= 14
    c.save()
    return buf.getvalue()


def _generate_pptx_bytes(text: str) -> bytes:
    """Generate a real .pptx file from text content."""
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    # Split by "Slide" markers or just add all text to one slide
    slides_text = text.split("Slide ")
    for i, slide_text in enumerate(slides_text):
        if not slide_text.strip():
            continue
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        body_shape = slide.placeholders[1]
        lines = slide_text.strip().split("\n")
        title_shape.text = lines[0][:60] if lines else f"Slide {i}"
        body_shape.text = "\n".join(lines[1:]) if len(lines) > 1 else ""
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _generate_xlsx_bytes(text: str) -> bytes:
    """Generate a real .xlsx file from text content with a multi-column table."""
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.title = "Recipe"
    for row_idx, line in enumerate(text.split("\n"), start=1):
        if line.strip():
            # Split on common delimiters to create multi-column layout
            parts = line.split(",") if "," in line else [line]
            for col_idx, part in enumerate(parts, start=1):
                ws.cell(row=row_idx, column=col_idx, value=part.strip())
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


FILE_GENERATORS = {
    "txt": _generate_txt_bytes,
    "csv": _generate_csv_bytes,
    "docx": _generate_docx_bytes,
    "pdf": _generate_pdf_bytes,
    "pptx": _generate_pptx_bytes,
    "xlsx": _generate_xlsx_bytes,
}


# ═══════════════════════════════════════════════════════════════════
# Seeder Class
# ═══════════════════════════════════════════════════════════════════

class DummyDataSeeder:
    """
    Orchestrates creation of all synthetic entities and ingests documents
    through the real production pipeline.
    """

    def __init__(self, db_url: str, vector_store_path: str, file_storage_path: str):
        self.db_url = db_url
        self.vector_store_path = vector_store_path
        self.file_storage_path = file_storage_path
        
        # Created entity IDs for test assertions
        self.role_ids: dict[str, str] = {}
        self.dept_ids: dict[str, str] = {}
        self.project_ids: dict[str, str] = {}
        self.user_ids: dict[str, str] = {}
        self.user_objects: dict[str, User] = {}
        self.document_ids: dict[str, str] = {}
        self.document_statuses: dict[str, str] = {}
        self.document_pii_flags: dict[str, bool] = {}
        self.total_chunks_indexed: int = 0

    async def seed_all(self) -> dict:
        """
        Runs the full seeding pipeline and returns a summary dict.
        """
        os.makedirs(self.vector_store_path, exist_ok=True)
        os.makedirs(self.file_storage_path, exist_ok=True)

        engine = create_async_engine(self.db_url, echo=False)

        async with engine.begin() as conn:
            await conn.run_sync(Base.metadata.create_all)

        session_factory = async_sessionmaker(engine, class_=AsyncSession, expire_on_commit=False)

        async with session_factory() as session:
            await self._seed_roles(session)
            await self._seed_departments(session)
            await self._seed_projects(session)
            await self._seed_users(session)
            await self._assign_project_memberships(session)
            await self._ingest_documents(session)

        await engine.dispose()

        return {
            "roles": len(self.role_ids),
            "departments": len(self.dept_ids),
            "projects": len(self.project_ids),
            "users": len(self.user_ids),
            "documents_total": len(self.document_ids),
            "documents_approved": sum(1 for s in self.document_statuses.values() if s == "approved"),
            "documents_quarantined": sum(1 for s in self.document_statuses.values() if s == "quarantined"),
            "total_chunks_indexed": self.total_chunks_indexed,
        }

    async def _seed_roles(self, session: AsyncSession):
        role_defs = [
            (Roles.PLATFORM_ADMIN, "Full platform admin access"),
            (Roles.DEPARTMENT_LEAD, "Department Lead access"),
            (Roles.PROJECT_ADMIN, "Project Admin access"),
            (Roles.EMPLOYEE, "Standard employee access"),
            (Roles.AUDITOR, "Compliance Auditor access"),
        ]
        for name, desc in role_defs:
            role = Role(name=name, description=desc)
            session.add(role)
            await session.flush()
            self.role_ids[name] = role.id
        await session.commit()

    async def _seed_departments(self, session: AsyncSession):
        dept_defs = [
            ("Quality Assurance", "QA", "Vehicle quality and testing"),
            ("Production", "PROD", "Manufacturing and assembly"),
            ("Human Resources", "HR", "Employee management"),
            ("Engineering", "ENG", "R&D and product engineering"),
            ("Supply Chain", "SCM", "Vendor and logistics management"),
        ]
        for name, code, desc in dept_defs:
            dept = Department(name=name, code=code, description=desc)
            session.add(dept)
            await session.flush()
            self.dept_ids[code] = dept.id
        await session.commit()

    async def _seed_projects(self, session: AsyncSession):
        project_defs = [
            ("proj_qa", "QA Welding Procedures", self.dept_ids["QA"]),
            ("proj_prod", "Production Supplier Audits", self.dept_ids["PROD"]),
            ("proj_hr", "HR Records Management", self.dept_ids["HR"]),
            ("proj_eng", "K15C Engine Development", self.dept_ids["ENG"]),
            ("proj_safety", "Plant Floor Safety", self.dept_ids["QA"]),
        ]
        for key, name, dept_id in project_defs:
            proj = Project(name=name, description=f"Project: {name}", department_id=dept_id)
            session.add(proj)
            await session.flush()
            self.project_ids[key] = proj.id
        await session.commit()

    async def _seed_users(self, session: AsyncSession):
        pwd_hash = bcrypt.hashpw("testpass123".encode("utf-8"), bcrypt.gensalt()).decode("utf-8")

        user_defs = [
            ("emp_user", "emp@marutisuzuki.com", self.role_ids[Roles.EMPLOYEE], self.dept_ids["QA"]),
            ("lead_user", "lead@marutisuzuki.com", self.role_ids[Roles.DEPARTMENT_LEAD], self.dept_ids["PROD"]),
            ("proj_admin_user", "projadmin@marutisuzuki.com", self.role_ids[Roles.PROJECT_ADMIN], self.dept_ids["ENG"]),
            ("plat_admin_user", "platadmin@marutisuzuki.com", self.role_ids[Roles.PLATFORM_ADMIN], None),
            ("auditor_user", "auditor@marutisuzuki.com", self.role_ids[Roles.AUDITOR], None),
            ("outsider_user", "outsider@marutisuzuki.com", self.role_ids[Roles.EMPLOYEE], self.dept_ids["HR"]),
        ]
        for username, email, role_id, dept_id in user_defs:
            user = User(
                username=username,
                email=email,
                hashed_password=pwd_hash,
                role_id=role_id,
                department_id=dept_id,
            )
            session.add(user)
            await session.flush()
            self.user_ids[username] = user.id
            # Reload role relationship
            await session.refresh(user, ["role"])
            self.user_objects[username] = user
        await session.commit()

    async def _assign_project_memberships(self, session: AsyncSession):
        """
        Assigns users to projects.
        Critical: outsider_user is deliberately NOT added to proj_qa.
        """
        memberships = [
            (self.user_ids["emp_user"], self.project_ids["proj_qa"]),
            (self.user_ids["emp_user"], self.project_ids["proj_safety"]),
            (self.user_ids["lead_user"], self.project_ids["proj_prod"]),
            (self.user_ids["proj_admin_user"], self.project_ids["proj_eng"]),
            (self.user_ids["plat_admin_user"], self.project_ids["proj_qa"]),
            (self.user_ids["plat_admin_user"], self.project_ids["proj_eng"]),
            (self.user_ids["auditor_user"], self.project_ids["proj_qa"]),
            # outsider_user intentionally NOT in proj_qa or proj_eng
            (self.user_ids["outsider_user"], self.project_ids["proj_hr"]),
        ]
        for user_id, project_id in memberships:
            await session.execute(
                insert(project_users).values(user_id=user_id, project_id=project_id)
            )
        await session.commit()

    async def _ingest_documents(self, session: AsyncSession):
        """
        Ingests all 6 synthetic documents through the real production pipeline.
        """
        from fastapi import UploadFile
        from services.document_service import ingest_document
        from adapters.vectorstore.faiss_adapter import vector_store
        from adapters.ai.embeddings import embedding_generator

        # Override vector store path for isolation
        vector_store.storage_dir = self.vector_store_path

        doc_project_map = {
            "sop_welding_robot.txt": ("proj_qa", "QA"),
            "supplier_audit_report.csv": ("proj_prod", "PROD"),
            "hr_employee_roster.docx": ("proj_hr", "HR"),
            "engine_spec_v2.pdf": ("proj_eng", "ENG"),
            "safety_training_deck.pptx": ("proj_qa", "QA"),
            "random_recipe.xlsx": ("proj_qa", "QA"),
        }

        for filename, doc_def in DOC_CONTENTS.items():
            proj_key, dept_code = doc_project_map[filename]
            project_id = self.project_ids[proj_key]
            department_id = self.dept_ids[dept_code]
            classification = doc_def["classification"]

            # Generate real binary file
            fmt = doc_def["format"]
            generator = FILE_GENERATORS[fmt]
            file_bytes = generator(doc_def["text"])

            # Create UploadFile compatible object
            file_obj = io.BytesIO(file_bytes)
            upload = UploadFile(filename=filename, file=file_obj)

            # Ingest through real pipeline
            doc = await ingest_document(
                db=session,
                upload_file=upload,
                project_id=project_id,
                department_id=department_id,
                user_provided_classification=classification,
                user_id=self.user_ids["plat_admin_user"],
            )

            self.document_ids[filename] = doc.id
            self.document_statuses[filename] = doc.status
            self.document_pii_flags[filename] = doc.pii_flagged

            # Count indexed chunks
            if doc.status == "approved" and doc.extracted_text:
                from rag.chunker import chunker
                chunks = chunker.split_text(doc.extracted_text)
                self.total_chunks_indexed += len(chunks)

            print(f"  [OK] {filename}: status={doc.status}, pii={doc.pii_flagged}, "
                  f"classification={doc.classification}")


async def run_seed(db_path: str = None) -> dict:
    """
    Convenience entry point for running the seeder.
    Returns the summary dict.
    """
    from core.config import get_settings
    settings = get_settings()

    db_url = settings.database_url
    vector_path = settings.vector_store_path
    storage_path = settings.file_storage_path

    # Extract db file path if SQLite to clean it up
    if db_url.startswith("sqlite"):
        db_file = db_url.split(":///")[1]
        if os.path.exists(db_file):
            print(f"Removing existing database file: {db_file}")
            os.remove(db_file)
    
    # Clean vector/storage dirs
    for d in (vector_path, storage_path):
        if os.path.exists(d):
            print(f"Cleaning directory: {d}")
            shutil.rmtree(d)
        os.makedirs(d, exist_ok=True)

    seeder = DummyDataSeeder(db_url, vector_path, storage_path)
    
    print(f"=== SEEDING DUMMY DATA TO ACTIVE DEV ENVIRONMENT ===")
    print(f"  DB URL: {db_url}")
    print(f"  Vector Store: {vector_path}")
    print(f"  File Storage: {storage_path}")
    print(f"====================================================")
    summary = await seeder.seed_all()
    
    print(f"\n=== SEED SUMMARY ===")
    print(f"  Roles: {summary['roles']}")
    print(f"  Departments: {summary['departments']}")
    print(f"  Projects: {summary['projects']}")
    print(f"  Users: {summary['users']}")
    print(f"  Documents: {summary['documents_total']} "
          f"(approved: {summary['documents_approved']}, quarantined: {summary['documents_quarantined']})")
    print(f"  Chunks indexed: {summary['total_chunks_indexed']}")
    print("=======================\n")

    return summary


if __name__ == "__main__":
    asyncio.run(run_seed())
